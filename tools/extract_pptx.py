"""Extract agent tables (text + hyperlinks) and images from the ABS agents deck."""
import json, os, sys, re
from pptx import Presentation
from pptx.util import Emu

SRC = sys.argv[1] if len(sys.argv) > 1 else r"C:\Users\cowi\Downloads\ABS Out-of-Box Agents - Customer ready.PPTX"
OUT = sys.argv[2] if len(sys.argv) > 2 else os.path.join(os.path.dirname(os.path.abspath(__file__)), "out")
os.makedirs(OUT, exist_ok=True)
os.makedirs(os.path.join(OUT, "media"), exist_ok=True)

prs = Presentation(SRC)


def cell_runs(cell):
    out = []
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            url = None
            try:
                if run.hyperlink and run.hyperlink.address:
                    url = run.hyperlink.address
            except Exception:
                pass
            out.append({"t": run.text, "u": url})
        out.append({"t": "\n", "u": None})
    return out


def runs_text(runs):
    return re.sub(r"\s+", " ", "".join(r["t"] for r in runs)).strip()


def runs_links(runs):
    links, cur, curu = [], "", None
    for r in runs:
        if r["u"] == curu:
            cur += r["t"]
        else:
            if curu:
                links.append({"label": cur.strip(), "url": curu})
            cur, curu = r["t"], r["u"]
    if curu:
        links.append({"label": cur.strip(), "url": curu})
    return [l for l in links if l["label"]]


def inch(v):
    return round(Emu(v).inches, 3) if v is not None else None


data = {"slides": []}
for si, slide in enumerate(prs.slides, 1):
    s = {"n": si, "tables": [], "texts": [], "pics": []}
    for shape in slide.shapes:
        if shape.has_table:
            rows = []
            for r in shape.table.rows:
                row = []
                for c in r.cells:
                    runs = cell_runs(c)
                    row.append({"text": runs_text(runs), "links": runs_links(runs)})
                rows.append(row)
            s["tables"].append({"top": inch(shape.top), "left": inch(shape.left), "rows": rows})
        elif shape.__class__.__name__ == "Picture":
            try:
                img = shape.image
                name = re.sub(r"[^A-Za-z0-9_.-]", "", f"s{si}_{shape.shape_id}_{shape.name}".replace(" ", "_"))
                name += "." + img.ext
                with open(os.path.join(OUT, "media", name), "wb") as f:
                    f.write(img.blob)
                s["pics"].append({"file": name, "top": inch(shape.top), "left": inch(shape.left),
                                  "w": inch(shape.width), "h": inch(shape.height)})
            except Exception as e:
                s["pics"].append({"err": str(e)})
        elif shape.has_text_frame:
            runs = []
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    u = None
                    try:
                        if r.hyperlink and r.hyperlink.address:
                            u = r.hyperlink.address
                    except Exception:
                        pass
                    runs.append({"t": r.text, "u": u})
                runs.append({"t": "\n", "u": None})
            t = runs_text(runs)
            if t:
                s["texts"].append({"text": t, "links": runs_links(runs),
                                   "top": inch(shape.top), "left": inch(shape.left)})
    data["slides"].append(s)

with open(os.path.join(OUT, "deck.json"), "w", encoding="utf-8") as f:
    json.dump(data, f, indent=1, ensure_ascii=False)
print("slides:", len(data["slides"]))
for s in data["slides"]:
    print(f"  slide {s['n']}: {len(s['tables'])} tables, {len(s['pics'])} pics, {len(s['texts'])} texts")
